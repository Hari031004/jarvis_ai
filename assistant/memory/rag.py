"""Document ingestion, local vector search, and citation generation."""

from __future__ import annotations

import csv
import hashlib
import html
import json
import math
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import numpy as np

from assistant.config import Settings
from assistant.memory.database import SQLiteStore, utc_now
from assistant.utils.logger import get_logger


@dataclass(slots=True)
class SearchResult:
    content: str
    citation: str
    score: float


class DocumentIntelligence:
    """Indexes documents into SQLite and searches them with hashed embeddings."""

    def __init__(self, settings: Settings, store: SQLiteStore) -> None:
        self.settings = settings
        self.store = store
        self.logger = get_logger(__name__)

    def ingest_path(self, path: Path) -> int:
        resolved = path.expanduser().resolve()
        if not resolved.exists():
            raise FileNotFoundError(resolved)
        text = self._extract_text(resolved)
        content_hash = hashlib.sha256(text.encode("utf-8")).hexdigest()
        existing = self.store.query_one("SELECT id, content_hash FROM documents WHERE path = ?", (str(resolved),))
        if existing and existing["content_hash"] == content_hash:
            return int(existing["id"])

        title = resolved.name
        now = utc_now()
        if existing:
            document_id = int(existing["id"])
            self.store.execute(
                "UPDATE documents SET title = ?, content_hash = ?, updated_at = ? WHERE id = ?",
                (title, content_hash, now, document_id),
            )
            self.store.execute("DELETE FROM document_chunks WHERE document_id = ?", (document_id,))
        else:
            document_id = self.store.insert_json(
                "documents",
                {
                    "path": str(resolved),
                    "title": title,
                    "content_hash": content_hash,
                    "metadata": {"suffix": resolved.suffix.lower()},
                    "created_at": now,
                    "updated_at": now,
                },
            )
        rows = []
        for index, chunk in enumerate(chunk_text(text, self.settings.rag_chunk_size, self.settings.rag_chunk_overlap)):
            rows.append(
                (
                    document_id,
                    index,
                    chunk,
                    json.dumps(embed_text(chunk, self.settings.rag_vector_dimensions)),
                    f"{resolved.name}#chunk-{index + 1}",
                )
            )
        self.store.executemany(
            "INSERT INTO document_chunks (document_id, chunk_index, content, vector, citation) VALUES (?, ?, ?, ?, ?)",
            rows,
        )
        return document_id

    def search(self, query: str, limit: int | None = None) -> list[SearchResult]:
        limit = limit or self.settings.rag_top_k
        query_vector = embed_text(query, self.settings.rag_vector_dimensions)
        rows = self.store.query("SELECT content, vector, citation FROM document_chunks")
        results = []
        for row in rows:
            vector = json.loads(row["vector"])
            score = cosine(query_vector, vector)
            if score > 0:
                results.append(SearchResult(row["content"], row["citation"], score))
        results.sort(key=lambda item: item.score, reverse=True)
        return results[:limit]

    def answer_context(self, query: str) -> str:
        results = self.search(query)
        if not results:
            return "No indexed document context matched the question."
        parts = []
        for result in results:
            parts.append(f"[{result.citation}] {result.content}")
        return "\n\n".join(parts)

    def _extract_text(self, path: Path) -> str:
        suffix = path.suffix.lower()
        if suffix in {".txt", ".md", ".py", ".js", ".ts", ".tsx", ".html", ".css", ".json", ".csv"}:
            return _read_text_family(path)
        if suffix == ".pdf":
            return _read_pdf(path)
        if suffix == ".docx":
            return _read_docx(path)
        if suffix in {".xlsx", ".xls"}:
            return _read_excel(path)
        if suffix == ".pptx":
            return _read_pptx(path)
        return path.read_text(encoding="utf-8", errors="ignore")


def chunk_text(text: str, size: int, overlap: int) -> Iterable[str]:
    cleaned = re.sub(r"\s+", " ", text).strip()
    if not cleaned:
        return []
    chunks = []
    start = 0
    size = max(200, size)
    overlap = max(0, min(overlap, size // 2))
    while start < len(cleaned):
        end = min(len(cleaned), start + size)
        chunks.append(cleaned[start:end])
        if end >= len(cleaned):
            break
        start = end - overlap
    return chunks


def embed_text(text: str, dimensions: int) -> list[float]:
    dimensions = max(32, dimensions)
    vector = np.zeros(dimensions, dtype=np.float32)
    for token in re.findall(r"[a-zA-Z0-9_]+", text.lower()):
        digest = hashlib.blake2b(token.encode("utf-8"), digest_size=8).digest()
        index = int.from_bytes(digest[:4], "little") % dimensions
        sign = 1.0 if digest[4] % 2 == 0 else -1.0
        vector[index] += sign
    norm = float(np.linalg.norm(vector))
    if norm:
        vector /= norm
    return [float(item) for item in vector]


def cosine(a: list[float], b: list[float]) -> float:
    left = np.array(a, dtype=np.float32)
    right = np.array(b, dtype=np.float32)
    denom = float(np.linalg.norm(left) * np.linalg.norm(right))
    return float(np.dot(left, right) / denom) if denom else 0.0


def _read_text_family(path: Path) -> str:
    if path.suffix.lower() == ".json":
        data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    if path.suffix.lower() == ".csv":
        with path.open("r", encoding="utf-8", errors="ignore", newline="") as file:
            rows = csv.reader(file)
            return "\n".join(" | ".join(row) for row in rows)
    if path.suffix.lower() == ".html":
        text = path.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r"<script.*?</script>|<style.*?</style>", " ", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<[^>]+>", " ", text)
        return html.unescape(text)
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _read_docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def _read_excel(path: Path) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(path, data_only=True, read_only=True)
    lines = []
    for sheet in workbook.worksheets:
        lines.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    lines = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"Slide {index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)
